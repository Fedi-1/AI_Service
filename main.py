import io
import json
import math
import os
import re
import traceback
import base64
import glob
import requests
from html import escape

import numpy as np
import pypdf
from PIL import Image, ImageDraw, ImageFont
from docx import Document as DocxDocument
from dotenv import load_dotenv
from fastapi import FastAPI, File, HTTPException, Query, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional
from groq import Groq
from pptx import Presentation

# â”€â”€â”€ Environment â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
if not GROQ_API_KEY:
    raise RuntimeError("GROQ_API_KEY is not set in the .env file")

ELEVENLABS_API_KEY = os.getenv("ELEVENLABS_API_KEY", "")

# ElevenLabs voice IDs
ELEVENLABS_VOICE_FR = "nPczCjzI2devNBz1zQrb"   # Brian â€” calm professional male (French)
ELEVENLABS_VOICE_EN = "21m00Tcm4TlvDq8ikWAM"   # Rachel â€” clear professional female (English)
ELEVENLABS_VOICE_AR = os.getenv("ELEVENLABS_VOICE_AR") or ELEVENLABS_VOICE_EN

client = Groq(api_key=GROQ_API_KEY)
MODEL  = "llama-3.3-70b-versatile"   # fast, free-tier Groq model

# â”€â”€â”€ App â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
app = FastAPI(title="LearnAI AI Service", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:8069", "http://localhost:4200"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# â”€â”€â”€ Constants â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
SUPPORTED_FILE_TYPES = {"PDF", "DOCX", "PPTX", "IMAGE"}
MIN_TEXT_LENGTH      = 50
CHUNK_WORD_SIZE      = 800
MAX_CHUNKS           = 10

os.makedirs("uploads/recap-videos", exist_ok=True)


def normalize_video_language(language: str | None) -> str:
    value = (language or "en").strip().lower()
    if value.startswith("fr"):
        return "fr"
    if value.startswith("ar") or "arab" in value or "عرب" in value:
        return "ar"
    return "en"


def video_language_name(language: str | None) -> str:
    code = normalize_video_language(language)
    if code == "fr":
        return "French"
    if code == "ar":
        return "Arabic"
    return "English"


# â”€â”€â”€ Text Extraction â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
def extract_text_pdf(file_bytes: bytes) -> str:
    reader = pypdf.PdfReader(io.BytesIO(file_bytes))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages).strip()


def extract_text_docx(file_bytes: bytes) -> str:
    doc = DocxDocument(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs).strip()


def extract_text_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines).strip()


def extract_text_image(file_bytes: bytes, filename: str = "") -> str:
    """Use Groq vision model for image text extraction."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    mime_map = {
        "jpg":  "image/jpeg",
        "jpeg": "image/jpeg",
        "png":  "image/png",
        "webp": "image/webp",
    }
    mime_type = mime_map.get(ext, "image/jpeg")
    b64_image = base64.b64encode(file_bytes).decode("utf-8")

    response = client.chat.completions.create(
        model="meta-llama/llama-4-scout-17b-16e-instruct",
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:{mime_type};base64,{b64_image}"},
                },
                {
                    "type": "text",
                    "text": (
                        "Extract all visible text from this image exactly as it appears. "
                        "Return only the extracted text with no commentary or formatting."
                    ),
                },
            ],
        }],
        max_tokens=4096,
    )
    content = response.choices[0].message.content.strip()
    if len(content) < 10:
        raise ValueError(
            "Vision model returned empty content. The image may not contain readable text "
            "or the model could not process it."
        )
    return content


def extract_text(file_bytes: bytes, file_type: str, filename: str = "") -> str:
    if file_type == "PDF":
        return extract_text_pdf(file_bytes)
    elif file_type == "DOCX":
        return extract_text_docx(file_bytes)
    elif file_type == "PPTX":
        return extract_text_pptx(file_bytes)
    elif file_type == "IMAGE":
        return extract_text_image(file_bytes, filename)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


# â”€â”€â”€ Text Chunking â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
def split_into_chunks(text: str) -> list[str]:
    words  = text.split()
    chunks = []
    for i in range(0, len(words), CHUNK_WORD_SIZE):
        chunk = " ".join(words[i: i + CHUNK_WORD_SIZE])
        chunks.append(chunk)
        if len(chunks) >= MAX_CHUNKS:
            break
    return chunks


# â”€â”€â”€ Groq Helpers â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
def clean_json_response(raw: str) -> str:
    """Strip markdown code fences and sanitize control characters inside JSON strings."""
    raw = raw.strip()
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)
    raw = raw.strip()

    # Replace literal newlines/tabs inside JSON string values with escaped versions
    # so that json.loads does not fail on invalid control characters.
    def replace_control_chars(m: re.Match) -> str:
        s = m.group(0)
        s = s.replace('\r\n', '\\n').replace('\r', '\\n').replace('\n', '\\n')
        s = s.replace('\t', '\\t')
        return s

    # Match everything between pairs of quotes (JSON string values)
    raw = re.sub(r'"(?:[^"\\]|\\.)*"', replace_control_chars, raw, flags=re.DOTALL)
    return raw


def groq_chat(prompt: str, temperature: float = 0.4, max_tokens: int = 8192) -> str:
    """Single chat completion call to Groq."""
    response = client.chat.completions.create(
        model=MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=temperature,
        max_tokens=max_tokens,
    )
    return response.choices[0].message.content.strip()


def generate_course_title(text_preview: str) -> str:
    prompt = (
        "Based on the following document excerpt, generate a short, clear, "
        "and descriptive course title (maximum 8 words). "
        "Return only the title text, nothing else.\n\n"
        f"Document excerpt:\n{text_preview[:600]}"
    )
    return groq_chat(prompt)


def generate_course_description(text_preview: str) -> str:
    prompt = (
        "Based on the following document excerpt, write a clear and informative course description "
        "in 2 to 3 sentences that explains what topics this course covers, what the student will learn, "
        "and why this knowledge is useful. "
        "Detect the language of the text and write the description in that same language. "
        "Return only the description text with no extra formatting or labels.\n\n"
        f"Document excerpt:\n{text_preview[:800]}"
    )
    return groq_chat(prompt)


def detect_category(text_preview: str) -> str:
    predefined = (
        "Informatique, MathÃ©matiques, Physique, Chimie, Biologie, "
        "Langues, Gestion, Droit, Ã‰lectronique, MÃ©canique"
    )
    prompt = (
        f"You are a document classifier. Read the following text excerpt and determine its academic or professional category.\n\n"
        f"First, try to match the document to one of these predefined categories:\n{predefined}\n\n"
        f"If the document clearly belongs to one of those categories, return that exact category name.\n"
        f"If the document does not match any of those predefined categories, invent the single most appropriate category name "
        f"based on the content (for example: Psychologie, Marketing, Architecture, Histoire, Philosophie, etc.).\n\n"
        f"Rules:\n"
        f"- Return ONLY the category name as plain text.\n"
        f"- No explanation, no punctuation around it, no extra words.\n"
        f"- Maximum 2 words.\n\n"
        f"Document excerpt:\n{text_preview[:500]}"
    )
    return groq_chat(prompt).strip().strip('."\'')


def generate_lesson(chunk: str, lesson_number: int, previous_titles: list[str]) -> dict:

    prior_context = ""
    if previous_titles:
        titles_list = "\n".join(f"  - Lesson {i+1}: {t}" for i, t in enumerate(previous_titles))
        prior_context = f"""
PREVIOUS LESSONS ALREADY TAUGHT (do NOT repeat or re-explain any concept or definition covered there):
{titles_list}

This lesson must build forward on those lessons. Assume the student already knows everything covered in the lessons above.
"""

    prompt = f"""You are a university professor creating a structured e-learning lesson. Your job is to TEACH, not to summarize.

Return ONLY a valid JSON object. No markdown, no code fences, no text outside the JSON.

LANGUAGE RULE: Detect the language of the text content below and write ALL fields (title, summary, content, questions, explanations, flashcards) in that SAME language. If the source is French: write in pure, correct French. Do NOT mix in Spanish, English, or any other language. If a technical term has no French equivalent, write that term in English but keep all surrounding text in French.
{prior_context}
FIELD DEFINITIONS â€” follow these strictly:

"summary":
- Write exactly 3 to 4 sentences in PAST TENSE, as a conclusion of what was just taught.
- It must reflect what the student now understands after completing this lesson.
- NEVER start with phrases like "this lesson covers", "we will explore", "students will learn", "ce cours couvre", "nous allons explorer", "les Ã©tudiants apprendront", or any forward-looking introduction.
- Start directly with what was explained â€” for example: "This lesson explored...", "We examined...", "Cette leÃ§on a mis en lumiÃ¨re...", "Nous avons vu que..."
- Do NOT teach or explain anything here. Only conclude.

"content":
- This is the main teaching section. Write it as a professor explaining to a student who has never seen this topic before.
- For EVERY concept in the text: explain WHAT it is, WHY it exists or matters, HOW it works in practice, and give a CONCRETE real-world example to illustrate it.
- Do NOT paraphrase or copy the source text. Teach the ideas from scratch using your own explanations.
- Do NOT re-explain or redefine any concept already covered in previous lessons listed above. If a concept was already taught, only reference it briefly and move forward.
- Structure the content in clear paragraphs with smooth transitions between ideas.
- The student must be able to fully understand and remember every concept after reading this field alone, without referring to the original document.
- Default minimum length: 5 full paragraphs, 600 words. EXCEPTION: if the source text chunk does not contain enough distinct concepts to fill 5 meaningful paragraphs, write 3 solid paragraphs of real content instead. Never repeat or rephrase the same idea multiple times just to reach a length target. Quality over quantity.
- The content field MUST be completely different from the summary field. Never copy the summary into content.

JSON structure:
{{
  "lessonNumber": {lesson_number},
  "title": "Descriptive lesson title",
  "summary": "Past-tense conclusion of what this lesson taught and what the student now understands.",
  "content": "Paragraph 1: Introduce the topic from scratch â€” what is it, where does it come from, why does it matter...\\n\\nParagraph 2: Explain the first key concept in depth â€” what it is, why it exists, how it works, with a real-world example...\\n\\nParagraph 3: Explain the second key concept in depth â€” what it is, why it exists, how it works, with a real-world example...\\n\\nParagraph 4: Explain the third key concept or go deeper into relationships between concepts, with comparisons or analogies...\\n\\nParagraph 5: Synthesize everything â€” how the concepts connect, what a student should take away, practical implications...",
  "estimatedReadTime": 8,
  "flashcards": [
    {{
      "term": "Key term",
      "definition": "Clear and concise definition.",
      "difficulty": "EASY | MEDIUM | HARD"
    }}
  ]
}}

Rules:
- "summary" MUST be in past tense â€” it is a conclusion, not an introduction
- "summary" MUST NOT start with forward-looking phrases (covers, will learn, allons explorer, etc.)
- "summary" and "content" MUST be completely different texts â€” never copy one into the other
- "content" MUST NOT re-explain concepts already covered in previous lessons
- "content" = professor-style teaching, every new concept explained with what/why/how/example; 5 paragraphs if the source has enough distinct concepts, otherwise 3 solid paragraphs â€” never pad with repeated ideas
- If source language is French: write in pure French only â€” no Spanish, no mixing; technical terms with no French equivalent may stay in English
- "estimatedReadTime" = estimated reading time in minutes for the content field (integer, minimum 1). Calculate it as: word count of the content field divided by 200, rounded up.
- For each flashcard "difficulty": assign EASY for basic definitions any beginner understands after one read; MEDIUM for concepts requiring understanding of relationships or how a mechanism works in practice; HARD for concepts requiring deep understanding, involving multiple interrelated ideas, or likely to be confused with similar concepts. You MUST distribute difficulty realistically â€” do not assign MEDIUM to all flashcards. A typical lesson must have a mix of EASY, MEDIUM, and HARD flashcards proportional to content complexity.
- Do NOT include literal newline characters inside JSON string values; use \\n for paragraph breaks
- Return ONLY the JSON object, nothing else

Text content to teach from:
{chunk}"""

    raw     = groq_chat(prompt)
    cleaned = clean_json_response(raw)
    lesson  = json.loads(cleaned)

    return lesson


# â”€â”€â”€ Recap Video â€” Script Generation â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
def generate_video_script(
    lesson_title: str,
    real_world_hook: str,
    key_takeaway: str,
    practical_tip: str,
    challenge_question: str,
    language: str,
) -> dict:
    """
    Generates narration scripts for all 4 slides via a single Llama 3.3 call.
    Returns dict with keys: slide1, slide2, slide3, slide4.
    """
    language_code = normalize_video_language(language)
    lang_name = video_language_name(language_code)
    language_detail = (
        "Use Modern Standard Arabic with clear educational phrasing."
        if language_code == "ar"
        else "Keep the language pure and natural."
    )
    prompt = f"""You are writing a narration script for a short educational video about "{lesson_title}". The video has 4 slides. Write the narrator text for each slide in {lang_name}. {language_detail} The narrator speaks in a calm, clear, professor-like tone. Use simple academic language appropriate for university students.

Slide 1 is about why this lesson matters in the real world. Base it on this hook: {real_world_hook}. Expand it into 3 sentences that explain the real-world importance naturally as if speaking to a student. Do not list facts â€” speak conversationally.

Slide 2 is about key concepts to remember. Base it on: {key_takeaway}. Write 3 to 4 sentences that introduce each key point conversationally, connecting them together naturally.

Slide 3 is about practical application. Base it on: {practical_tip}. Write EXACTLY 3 separate practical steps. Format them strictly as 3 lines: line 1 starts with '1. ' followed by the step text, line 2 starts with '2. ' followed by the step text, line 3 starts with '3. ' followed by the step text. Each step is one complete sentence describing a specific concrete action using a real tool or technology. Do NOT combine steps into one paragraph. Do NOT add any text before or after the 3 numbered lines. Return only the 3 lines.

Slide 4 is a challenge question to prepare the student for the quiz. Use this question: {challenge_question}. Write 2 sentences: first tell the student to think carefully about the question, then state the question clearly.

Return ONLY a valid JSON object with exactly 4 fields: slide1, slide2, slide3, slide4. Each field contains the narrator text as a plain string with no formatting."""

    try:
        raw     = groq_chat(prompt, temperature=0.6, max_tokens=600)
        cleaned = clean_json_response(raw)
        scripts = json.loads(cleaned)
        for key in ("slide1", "slide2", "slide3", "slide4"):
            if key not in scripts or not isinstance(scripts[key], str):
                raise ValueError(f"Missing or invalid key: {key}")
        return scripts
    except Exception as e:
        print(f"[generate_video_script] fallback due to: {e}")
        if language_code == "fr":
            return {
                "slide1": f"{real_world_hook} Ce concept joue un rÃ´le fondamental dans de nombreuses technologies modernes que vous utilisez au quotidien.",
                "slide2": f"{key_takeaway} Ces points essentiels forment la base de votre comprÃ©hension de ce sujet.",
                "slide3": f"1. {practical_tip}\n2. Pratiquez avec un exemple concret sur votre machine.\n3. Comparez votre rÃ©sultat avec la documentation officielle.",
                "slide4": f"Prenez un moment pour rÃ©flÃ©chir attentivement Ã  cette question avant de passer au quiz. {challenge_question}",
            }
        if language_code == "ar":
            return {
                "slide1": f"{real_world_hook} هذا المفهوم مهم في التقنيات الحديثة التي نستخدمها يوميا.",
                "slide2": f"{key_takeaway} هذه النقاط الأساسية تساعدك على فهم الموضوع بوضوح.",
                "slide3": f"1. {practical_tip}\n2. طبق الفكرة على مثال بسيط في بيئة العمل.\n3. قارن النتيجة مع مرجع موثوق.",
                "slide4": f"خذ لحظة للتفكير في هذا السؤال قبل الانتقال إلى الاختبار. {challenge_question}",
            }
        else:
            return {
                "slide1": f"{real_world_hook} This concept plays a fundamental role in many modern technologies you use every day.",
                "slide2": f"{key_takeaway} These essential points form the foundation of your understanding of this topic.",
                "slide3": f"1. {practical_tip}\n2. Practice with a concrete example in your development environment.\n3. Compare your result against the official documentation.",
                "slide4": f"Take a moment to think carefully about this question before moving on to the quiz. {challenge_question}",
            }


# â”€â”€â”€ Recap Video â€” ElevenLabs Audio Generation â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
def generate_slide_audio(
    slide_text: str,
    language: str,
    slide_index: int,
) -> tuple[str, list[dict]]:
    """
    Calls ElevenLabs TTS with word-level timestamps.
    Returns (mp3_path, word_timestamps_list).
    Falls back to silent audio on any error.
    """
    os.makedirs(os.path.join("uploads", "recap-videos"), exist_ok=True)
    mp3_path = os.path.abspath(
        os.path.join("uploads", "recap-videos", f"temp_slide_{slide_index}.mp3")
    )

    def _silent_fallback() -> tuple[str, list]:
        """Write 20 s of silent MP3 using a minimal WAVâ†’MP3 trick via numpy."""
        # Write a 20-second silent WAV, convert via moviepy/ffmpeg
        try:
            import wave, struct
            silent_wav = mp3_path.replace(".mp3", "_silent.wav")
            sample_rate = 22050
            num_samples = sample_rate * 20
            with wave.open(silent_wav, "w") as wf:
                wf.setnchannels(1)
                wf.setsampwidth(2)
                wf.setframerate(sample_rate)
                wf.writeframes(struct.pack("<" + "h" * num_samples, *([0] * num_samples)))
            # Convert WAV â†’ MP3 with ffmpeg (moviepy bundles it)
            from moviepy import AudioFileClip as AFC
            afc = AFC(silent_wav)
            afc.write_audiofile(mp3_path, logger=None)
            afc.close()
            os.remove(silent_wav)
        except Exception as ex:
            print(f"[generate_slide_audio] silent fallback failed: {ex}")
            # Last resort: write empty bytes so MoviePy can still open it
            open(mp3_path, "wb").close()
        return mp3_path, []

    if not ELEVENLABS_API_KEY:
        print(f"[generate_slide_audio] ELEVENLABS_API_KEY not set â€” using silent fallback")
        return _silent_fallback()

    language_code = normalize_video_language(language)
    if language_code == "fr":
        voice_id = ELEVENLABS_VOICE_FR
    elif language_code == "ar":
        voice_id = ELEVENLABS_VOICE_AR
    else:
        voice_id = ELEVENLABS_VOICE_EN
    url = f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}/with-timestamps"

    headers = {
        "xi-api-key": ELEVENLABS_API_KEY,
        "Content-Type": "application/json",
    }
    body = {
        "model_id": "eleven_multilingual_v2",
        "text": slide_text,
        "voice_settings": {
            "stability": 0.75,
            "similarity_boost": 0.85,
            "style": 0.3,
            "use_speaker_boost": True,
        },
    }

    try:
        resp = requests.post(url, headers=headers, json=body, timeout=60)
        resp.raise_for_status()
        data = resp.json()

        # Decode audio
        audio_bytes = base64.b64decode(data["audio_base64"])
        with open(mp3_path, "wb") as f:
            f.write(audio_bytes)

        # Parse word-level timestamps from character alignment
        alignment  = data.get("alignment", {})
        chars      = alignment.get("characters", [])
        starts     = alignment.get("character_start_times_seconds", [])
        ends       = alignment.get("character_end_times_seconds", [])

        word_timestamps: list[dict] = []
        current_word = ""
        word_start   = 0.0
        word_end     = 0.0

        for i, ch in enumerate(chars):
            char_start = starts[i] if i < len(starts) else 0.0
            char_end   = ends[i]   if i < len(ends)   else 0.0
            if ch == " " or ch == "\n":
                if current_word:
                    word_timestamps.append({
                        "word":  current_word,
                        "start": word_start,
                        "end":   word_end,
                    })
                    current_word = ""
            else:
                if not current_word:
                    word_start = char_start
                current_word += ch
                word_end = char_end
        if current_word:
            word_timestamps.append({
                "word":  current_word,
                "start": word_start,
                "end":   word_end,
            })

        print(f"[generate_slide_audio] slide {slide_index}: {len(word_timestamps)} words, "
              f"audio {os.path.getsize(mp3_path):,} bytes")
        return mp3_path, word_timestamps

    except Exception as e:
        print(f"[generate_slide_audio] ElevenLabs error for slide {slide_index}: {e}")
        traceback.print_exc()
        return _silent_fallback()


# â”€â”€â”€ Recap Video â€” Audio Duration Helper â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
def _audio_duration_seconds(mp3_path: str) -> float:
    """Return duration of an MP3 file in seconds using mutagen, or estimate from file size."""
    try:
        from mutagen.mp3 import MP3
        return MP3(mp3_path).info.length
    except Exception:
        pass
    try:
        size = os.path.getsize(mp3_path)
        # Rough estimate: 128 kbps â†’ ~16000 bytes/sec
        return max(1.0, size / 16000)
    except Exception:
        return 10.0


def _clean_caption_word(word: object) -> str:
    return re.sub(r"\s+", " ", str(word or "")).strip()


def _caption_text_from_words(words: list[dict]) -> str:
    return " ".join(
        cleaned for cleaned in (_clean_caption_word(w.get("word", "")) for w in words) if cleaned
    )


def _subtitle_cues_from_words(words: list[dict], language: str) -> list[dict]:
    cleaned_words = []
    for word in words:
        try:
            text = _clean_caption_word(word.get("word", ""))
            start = float(word.get("start", 0))
            end = float(word.get("end", start))
        except (TypeError, ValueError):
            continue
        if text:
            cleaned_words.append({"word": text, "start": start, "end": end})

    if not cleaned_words:
        return []

    language_code = normalize_video_language(language)
    max_words = 8 if language_code == "ar" else 10
    max_chars = 46 if language_code == "ar" else 58
    min_words_before_soft_break = 3 if language_code == "ar" else 4
    raw_cues: list[dict] = []
    cue_words: list[dict] = []
    cue_start = cleaned_words[0]["start"]

    for i, word in enumerate(cleaned_words):
        next_word = cleaned_words[i + 1] if i + 1 < len(cleaned_words) else None
        if not cue_words:
            cue_start = word["start"]
        cue_words.append(word)

        text = _caption_text_from_words(cue_words)
        enough_for_soft_break = len(cue_words) >= min_words_before_soft_break
        ends_sentence = bool(re.search(r"[.!?\u061f\u2026]$", word["word"]))
        ends_phrase = bool(re.search(r"[,;:\u060c\u061b]$", word["word"]))
        next_gap = max(0.0, next_word["start"] - word["end"]) if next_word else 0.0
        should_break = (
            i == len(cleaned_words) - 1
            or (ends_sentence and enough_for_soft_break)
            or next_gap >= 0.45
            or len(cue_words) >= max_words
            or len(text) >= max_chars
            or (ends_phrase and len(cue_words) >= 5)
        )

        if should_break:
            raw_cues.append({
                "text": text,
                "start": cue_start,
                "end": max(word["end"], cue_start + 0.25),
            })
            cue_words = []

    cues = []
    for index, cue in enumerate(raw_cues):
        next_cue = raw_cues[index + 1] if index + 1 < len(raw_cues) else None
        min_end = cue["start"] + 1.05
        padded_end = cue["end"] + 0.45
        if next_cue:
            next_limit = next_cue["start"] - 0.04
            end = max(cue["end"], min(next_limit, max(min_end, padded_end)))
        else:
            end = max(min_end, padded_end)
        cues.append({**cue, "end": end})
    return cues


def _subtitle_cues_from_script(script: str, duration_seconds: float, language: str) -> list[dict]:
    words = re.sub(r"\s+", " ", script or "").strip().split()
    if not words:
        return []
    duration = max(1.0, float(duration_seconds or 1.0))
    synthetic_words = [
        {
            "word": word,
            "start": (index / len(words)) * duration,
            "end": ((index + 1) / len(words)) * duration,
        }
        for index, word in enumerate(words)
    ]
    return _subtitle_cues_from_words(synthetic_words, language)


def _full_video_subtitle_cues(
    slide_configs: list[dict],
    slide_audio_data: list[dict],
    language: str,
) -> list[dict]:
    full_cues: list[dict] = []
    offset = 0.0
    for config, audio_data in zip(slide_configs, slide_audio_data):
        duration = float(audio_data.get("audioDurationSeconds") or 0)
        words = audio_data.get("words") or []
        cues = (
            _subtitle_cues_from_words(words, language)
            if words
            else _subtitle_cues_from_script(config.get("script", ""), duration, language)
        )
        for cue in cues:
            full_cues.append({
                "text": cue["text"],
                "start": cue["start"] + offset,
                "end": cue["end"] + offset,
            })
        offset += math.ceil(duration * 30) / 30
    return full_cues


def _translate_subtitle_cues(
    cues: list[dict],
    target_language: str,
    source_language: str,
) -> list[dict]:
    target_code = normalize_video_language(target_language)
    source_code = normalize_video_language(source_language)
    if not cues or target_code == source_code:
        return [dict(cue) for cue in cues]

    translated_texts: list[str] = []
    batch_size = 10
    for start in range(0, len(cues), batch_size):
        batch = cues[start:start + batch_size]
        batch_texts = [cue["text"] for cue in batch]
        translated_texts.extend(
            _translate_subtitle_batch(batch_texts, target_code, source_code)
        )

    if len(translated_texts) != len(cues):
        print(f"[translate_subtitle_cues] length mismatch for {target_code}; keeping source")
        return [dict(cue) for cue in cues]

    if _looks_untranslated([cue["text"] for cue in cues], translated_texts, target_code):
        print(f"[translate_subtitle_cues] detected untranslated {target_code} batch; retrying one by one")
        translated_texts = [
            _translate_single_subtitle(cue["text"], target_code, source_code)
            for cue in cues
        ]

    return [
        {**cue, "text": _clean_caption_word(translated_texts[index]) or cue["text"]}
        for index, cue in enumerate(cues)
    ]


def _translate_subtitle_batch(
    texts: list[str],
    target_language: str,
    source_language: str,
) -> list[str]:
    prompt = f"""Translate these subtitle cues from {video_language_name(source_language)} to {video_language_name(target_language)}.

Return ONLY this JSON shape:
{{"items":["translated cue 1","translated cue 2"]}}

Rules:
- The "items" array must contain exactly {len(texts)} strings.
- Keep each cue short and natural for subtitles.
- Preserve the meaning and order.
- Do not copy the source language unless the item is a proper noun or technical term.
- No markdown, no numbering, no explanations.

Source cues:
{json.dumps(texts, ensure_ascii=False)}"""

    try:
        raw = groq_chat(prompt, temperature=0.1, max_tokens=max(900, len(texts) * 120))
        parsed = _parse_translation_response(raw)
        if len(parsed) != len(texts):
            raise ValueError(f"expected {len(texts)} translations, got {len(parsed)}")
        if _looks_untranslated(texts, parsed, target_language):
            raise ValueError("batch appears untranslated")
        return parsed
    except Exception as e:
        print(f"[translate_subtitle_batch] retrying individually for {target_language}: {e}")
        return [
            _translate_single_subtitle(text, target_language, source_language)
            for text in texts
        ]


def _parse_translation_response(raw: str) -> list[str]:
    cleaned = clean_json_response(raw)
    try:
        parsed = json.loads(cleaned)
        if isinstance(parsed, list):
            return [_clean_caption_word(item) for item in parsed]
        if isinstance(parsed, dict) and isinstance(parsed.get("items"), list):
            return [_clean_caption_word(item) for item in parsed["items"]]
    except Exception:
        pass

    lines = []
    for line in cleaned.splitlines():
        item = re.sub(r"^\s*(?:[-*]|\d+[.)])\s*", "", line).strip().strip('"')
        if item:
            lines.append(_clean_caption_word(item))
    return lines


def _translate_single_subtitle(text: str, target_language: str, source_language: str) -> str:
    prompt = (
        f"Translate this subtitle from {video_language_name(source_language)} "
        f"to {video_language_name(target_language)}. Return only the translated subtitle, "
        f"with no quotes and no explanation:\n{text}"
    )
    try:
        translated = _clean_caption_word(groq_chat(prompt, temperature=0.1, max_tokens=180))
        translated = translated.strip().strip('"')
        if translated and not _looks_untranslated([text], [translated], target_language):
            return translated
    except Exception as e:
        print(f"[translate_single_subtitle] keeping source for {target_language}: {e}")
    return text


def _looks_untranslated(source_texts: list[str], translated_texts: list[str], target_language: str) -> bool:
    if not source_texts or not translated_texts:
        return True

    normalized_source = [_clean_caption_word(text).lower() for text in source_texts]
    normalized_translated = [_clean_caption_word(text).lower() for text in translated_texts]
    identical = sum(
        1 for source, translated in zip(normalized_source, normalized_translated)
        if source and source == translated
    )

    if identical >= max(1, int(len(source_texts) * 0.6)):
        return True

    combined = " ".join(translated_texts)
    if target_language == "ar":
        return not bool(re.search(r"[\u0600-\u06ff]", combined))

    if target_language == "en":
        french_markers = re.findall(
            r"\b(le|la|les|un|une|des|et|est|en|de|du|pour|avec|sur|dans|qui|que|ce|cette|ces|nous|vous|ils|elles)\b",
            combined.lower(),
        )
        return len(french_markers) >= max(4, len(translated_texts) // 3)

    return False


def _format_vtt_timestamp(seconds: float) -> str:
    total_ms = max(0, int(round(seconds * 1000)))
    hours, remainder = divmod(total_ms, 3_600_000)
    minutes, remainder = divmod(remainder, 60_000)
    secs, millis = divmod(remainder, 1000)
    return f"{hours:02}:{minutes:02}:{secs:02}.{millis:03}"


def _write_vtt_file(path: str, cues: list[dict], language: str) -> None:
    os.makedirs(os.path.dirname(path), exist_ok=True)
    language_code = normalize_video_language(language)
    lines = ["WEBVTT", ""]

    for index, cue in enumerate(cues, start=1):
        text = _clean_caption_word(cue.get("text", ""))
        if not text:
            continue
        if language_code == "ar":
            text = f"\u202b{text}\u202c"
        start = float(cue.get("start", 0))
        end = max(start + 0.25, float(cue.get("end", start + 0.25)))
        lines.extend([
            str(index),
            f"{_format_vtt_timestamp(start)} --> {_format_vtt_timestamp(end)} line:86% position:50% align:center size:86%",
            escape(text, quote=False),
            "",
        ])

    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))


def _subtitle_paths_for_video(video_path: str) -> dict[str, str]:
    base, _ = os.path.splitext(video_path)
    return {
        "subtitleEnPath": f"{base}.en.vtt",
        "subtitleFrPath": f"{base}.fr.vtt",
        "subtitleArPath": f"{base}.ar.vtt",
    }


def _recap_response(video_path: str) -> dict[str, str]:
    response = {"recapVideoPath": video_path}
    for key, subtitle_path in _subtitle_paths_for_video(video_path).items():
        if os.path.exists(subtitle_path):
            response[key] = subtitle_path
    return response


def _write_subtitle_tracks(video_path: str, source_cues: list[dict], source_language: str) -> dict[str, str]:
    paths = _subtitle_paths_for_video(video_path)
    for language_code, key in (("en", "subtitleEnPath"), ("fr", "subtitleFrPath"), ("ar", "subtitleArPath")):
        translated_cues = _translate_subtitle_cues(source_cues, language_code, source_language)
        _write_vtt_file(paths[key], translated_cues, language_code)
    return paths


# â”€â”€â”€ Recap Video Generation â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
def generate_recap_video(
    lesson_number: int,
    lesson_title: str,
    flashcard_terms: list[str],
    lesson_summary: str,
    estimated_read_time: int,
    course_title: str,
    language: str,
) -> dict[str, str] | None:
    """
    Render a narrated MP4 recap video and WebVTT subtitle tracks.
    Returns paths to the video and subtitle files, or None on failure.
    """
    print(f"[generate_recap_video] called for lesson {lesson_number}: '{lesson_title}'")
    try:
        language_code = normalize_video_language(language)
        is_fr = language_code == "fr"
        is_ar = language_code == "ar"

        # â”€â”€ Content generation (Groq) â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        try:
            if is_fr:
                hook_prompt = (
                    f"Donne 2 Ã  3 phrases expliquant des applications du monde rÃ©el oÃ¹ \"{lesson_title}\" est utilisÃ©. "
                    f"Pour chaque application mentionne une app ou technologie bien connue et explique comment ce concept s'y applique. "
                    f"Maximum 50 mots au total. RÃ©ponds en franÃ§ais."
                )
            elif is_ar:
                hook_prompt = (
                    f"Give 2 to 3 sentences explaining real-world applications where \"{lesson_title}\" is used. "
                    f"For each application mention a specific well-known app or technology and explain how this concept applies to it. "
                    f"Total maximum 50 words. Respond in Modern Standard Arabic."
                )
            else:
                hook_prompt = (
                    f"Give 2 to 3 sentences explaining real-world applications where \"{lesson_title}\" is used. "
                    f"For each application mention a specific well-known app or technology and explain how this concept applies to it. "
                    f"Total maximum 50 words. Respond in English."
                )
            real_world_hook = groq_chat(hook_prompt, temperature=0.7, max_tokens=200).strip()
        except Exception:
            if is_fr:
                real_world_hook = "Ce concept est utilisÃ© dans de nombreuses applications modernes comme les bases de donnÃ©es et les frameworks web."
            elif is_ar:
                real_world_hook = "يستخدم هذا المفهوم في تطبيقات حديثة مثل قواعد البيانات وأطر عمل الويب."
            else:
                real_world_hook = "This concept is used in modern applications like databases and web frameworks."

        try:
            if is_fr:
                takeaway_prompt = (
                    f"Liste exactement 3 points clÃ©s qu'un Ã©tudiant doit retenir d'une leÃ§on sur \"{lesson_title}\". "
                    f"Format : 3 courtes affirmations commenÃ§ant chacune par un tiret. Maximum 12 mots par point. RÃ©ponds en franÃ§ais."
                )
            elif is_ar:
                takeaway_prompt = (
                    f"List exactly 3 key points a student must remember from a lesson about \"{lesson_title}\". "
                    f"Format as 3 short statements each starting with a dash character. Each point maximum 12 words. Respond in Modern Standard Arabic."
                )
            else:
                takeaway_prompt = (
                    f"List exactly 3 key points a student must remember from a lesson about \"{lesson_title}\". "
                    f"Format as 3 short statements each starting with a dash character. Each point maximum 12 words. Respond in English."
                )
            key_takeaway = groq_chat(takeaway_prompt, temperature=0.7, max_tokens=200).strip()
        except Exception:
            if is_fr:
                key_takeaway = "- MaÃ®trisez les concepts fondamentaux avant d'avancer.\n- Pratiquez rÃ©guliÃ¨rement avec des exemples concrets.\n- Reliez ce concept aux notions dÃ©jÃ  apprises."
            elif is_ar:
                key_takeaway = "- افهم المفاهيم الأساسية قبل التقدم.\n- تدرب بانتظام على أمثلة واقعية.\n- اربط الفكرة بما تعلمته سابقا."
            else:
                key_takeaway = "- Master the fundamentals before moving forward.\n- Practice regularly with concrete examples.\n- Connect this concept to what you already know."

        try:
            if is_fr:
                tip_prompt = (
                    f"Donne exactement 3 Ã©tapes pratiques concrÃ¨tes pour quelqu'un qui apprend \"{lesson_title}\". "
                    f"Chaque Ã©tape doit mentionner un outil, une mÃ©thode ou une technologie spÃ©cifique. "
                    f"Retourne exactement 3 lignes formatÃ©es strictement comme '1. [Ã©tape]' retour Ã  la ligne '2. [Ã©tape]' retour Ã  la ligne '3. [Ã©tape]'. "
                    f"Maximum 18 mots par Ã©tape. RÃ©ponds en franÃ§ais. Retourne uniquement les 3 lignes sans texte supplÃ©mentaire."
                )
            elif is_ar:
                tip_prompt = (
                    f"Give exactly 3 concrete practical steps for someone learning about \"{lesson_title}\". "
                    f"Each step must mention a specific tool, method, or technology. "
                    f"Return exactly 3 lines formatted strictly as '1. [step]' newline '2. [step]' newline '3. [step]'. "
                    f"Maximum 18 words per step. Respond in Modern Standard Arabic. Return only the 3 lines with no extra text or explanation."
                )
            else:
                tip_prompt = (
                    f"Give exactly 3 concrete practical steps for someone learning about \"{lesson_title}\". "
                    f"Each step must mention a specific tool, method, or technology. "
                    f"Return exactly 3 lines formatted strictly as '1. [step]' newline '2. [step]' newline '3. [step]'. "
                    f"Maximum 18 words per step. Respond in English. Return only the 3 lines with no extra text or explanation."
                )
            practical_tip = groq_chat(tip_prompt, temperature=0.7, max_tokens=200).strip()
        except Exception:
            if is_fr:
                practical_tip = "1. Lisez le rÃ©sumÃ© et identifiez les concepts clÃ©s.\n2. CrÃ©ez un exemple simple.\n3. Expliquez le concept Ã  voix haute."
            elif is_ar:
                practical_tip = "1. اقرأ الملخص وحدد المفاهيم الأساسية.\n2. ابن مثالا بسيطا لتطبيق الفكرة.\n3. اشرح المفهوم بصوت واضح."
            else:
                practical_tip = "1. Read the summary and identify the key concepts.\n2. Build a simple example.\n3. Explain the concept out loud."

        try:
            if is_fr:
                challenge_prompt = (
                    f"Formule une question de rÃ©flexion ouverte sur \"{lesson_title}\" pour prÃ©parer un Ã©tudiant Ã  un quiz. "
                    f"La question doit tester la comprÃ©hension profonde, pas la mÃ©morisation. Maximum 25 mots. RÃ©ponds en franÃ§ais."
                )
            elif is_ar:
                challenge_prompt = (
                    f"Write one open-ended thinking question about \"{lesson_title}\" to prepare a student for a quiz. "
                    f"The question should test deep understanding, not memorisation. Maximum 25 words. Respond in Modern Standard Arabic."
                )
            else:
                challenge_prompt = (
                    f"Write one open-ended thinking question about \"{lesson_title}\" to prepare a student for a quiz. "
                    f"The question should test deep understanding, not memorisation. Maximum 25 words. Respond in English."
                )
            challenge_question = groq_chat(challenge_prompt, temperature=0.7, max_tokens=100).strip()
        except Exception:
            if is_fr:
                challenge_question = f"Comment appliqueriez-vous les concepts de \"{lesson_title}\" pour rÃ©soudre un problÃ¨me concret ?"
            elif is_ar:
                challenge_question = f"كيف تطبق مفاهيم \"{lesson_title}\" لحل مشكلة واقعية؟"
            else:
                challenge_question = f"How would you apply the concepts of \"{lesson_title}\" to solve a real-world problem?"

        # â”€â”€ Narration scripts (Groq) â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        scripts = generate_video_script(
            lesson_title=lesson_title,
            real_world_hook=real_world_hook,
            key_takeaway=key_takeaway,
            practical_tip=practical_tip,
            challenge_question=challenge_question,
            language=language_code,
        )

        # â”€â”€ Slide configurations â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        GOLD_HEX   = "#f59e0b"
        INDIGO_HEX = "#6366f1"

        slide_configs = [
            {
                "title":      "Pourquoi apprendre ceci ?" if is_fr else "لماذا هذا مهم؟" if is_ar else "Why does this matter?",
                "accentColor": GOLD_HEX,
                "script":     scripts["slide1"],
            },
            {
                "title":      "Ã€ retenir" if is_fr else "أهم الأفكار" if is_ar else "Key Takeaway",
                "accentColor": INDIGO_HEX,
                "script":     scripts["slide2"],
            },
            {
                "title":      "Dans la pratique" if is_fr else "في التطبيق" if is_ar else "In Practice",
                "accentColor": GOLD_HEX,
                "script":     scripts["slide3"],
            },
            {
                "title":      "DÃ©fi avant le quiz" if is_fr else "تحد قبل الاختبار" if is_ar else "Challenge before the quiz",
                "accentColor": "#a78bfa",
                "script":     scripts["slide4"],
            },
        ]

        # â”€â”€ Generate audio for all 4 slides â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        slide_audio_data = []
        for i, cfg in enumerate(slide_configs):
            slide_num = i + 1
            print(f"[generate_recap_video] generating audio for slide {slide_num}...")
            audio_path, word_ts = generate_slide_audio(
                slide_text=cfg["script"],
                language=language_code,
                slide_index=slide_num,
            )
            duration = _audio_duration_seconds(audio_path)
            slide_audio_data.append({
                "audioFilePath": audio_path,
                "words":         word_ts,
                "audioDurationSeconds": duration,
            })

        # â”€â”€ Build VideoData JSON for Remotion â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        video_data = {
            "lessonTitle":     lesson_title,
            "language":        language_code,
            "flashcardCount":  len(flashcard_terms),
            "quizCount":       5,
            "estimatedReadTime": estimated_read_time,
            "slides": [
                {
                    "title":               slide_configs[i]["title"],
                    "accentColor":         slide_configs[i]["accentColor"],
                    "script":              slide_configs[i]["script"],
                    "words":               slide_audio_data[i]["words"],
                    "showSubtitles":       False,
                    "audioDurationSeconds": slide_audio_data[i]["audioDurationSeconds"],
                    "audioFilePath":       slide_audio_data[i]["audioFilePath"],
                }
                for i in range(4)
            ],
        }

        safe = re.sub(r"[^\w]", "_", lesson_title.encode("ascii", "ignore").decode())[:30]
        os.makedirs(os.path.join("uploads", "recap-videos"), exist_ok=True)

        json_path = os.path.abspath(
            os.path.join("uploads", "recap-videos", f"data_{lesson_number}_{language_code}_{safe}.json")
        )
        out = os.path.abspath(
            os.path.join("uploads", "recap-videos", f"lesson_{lesson_number}_{language_code}_{safe}.mp4")
        )

        with open(json_path, "w", encoding="utf-8") as jf:
            json.dump(video_data, jf, ensure_ascii=False, indent=2)

        # â”€â”€ Call Remotion renderer â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        import subprocess
        project_dir = os.path.abspath(os.path.dirname(__file__))
        render_script = os.path.join(project_dir, "remotion-renderer", "render.mjs")

        print(f"[generate_recap_video] calling Remotion renderer â†’ {out}")
        result = subprocess.run(
            ["node", render_script, "--data", json_path, "--output", out],
            cwd=project_dir,
            capture_output=True,
            text=True,
            timeout=600,
        )

        if result.returncode != 0:
            print(f"[generate_recap_video] Remotion stderr:\n{result.stderr}")
            print(f"[generate_recap_video] Remotion stdout:\n{result.stdout}")
            return None

        print(result.stdout)
        source_cues = _full_video_subtitle_cues(slide_configs, slide_audio_data, language_code)
        _write_subtitle_tracks(out, source_cues, language_code)

        # â”€â”€ Cleanup â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        try:
            os.remove(json_path)
        except Exception:
            pass
        for tmp in glob.glob(os.path.join("uploads", "recap-videos", "temp_slide_*.mp3")):
            try:
                os.remove(tmp)
            except Exception:
                pass

        print(f"[generate_recap_video] done â€” {os.path.getsize(out):,} bytes")
        return _recap_response(out)

    except Exception as exc:
        print(f"[generate_recap_video] ERROR: {exc}")
        traceback.print_exc()
        return None



# â”€â”€â”€ On-demand Quiz Generation â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
def generate_quiz_questions(lesson_content: str, previous_questions: list[str] = []) -> list:
    """Generate exactly 5 fresh quiz questions from lesson content using high-temperature sampling."""
    exclusion_block = ""
    if previous_questions:
        formatted = "\n".join(f"  - {q}" for q in previous_questions)
        exclusion_block = f"""
FORBIDDEN QUESTIONS â€” these questions were already used in previous attempts. You MUST NOT generate any question that is similar in meaning, structure, or topic to any of the following:
{formatted}

If you repeat or closely paraphrase any of these, it is a critical failure.
"""

    prompt = f"""You are a quiz generator for an e-learning platform. Your task is to generate exactly 5 quiz questions based on the lesson content provided.
{exclusion_block}
CRITICAL INSTRUCTION â€” ANGLE OF APPROACH:
Every time you are called you must approach the lesson from a completely different angle.
Do NOT generate questions about the most obvious or prominent concepts â€” instead pick less obvious facts, edge cases, nuances, or relationships between concepts that a student might overlook or take for granted.
Force the student to think deeply, not just recall headline facts.

LANGUAGE RULE: Detect the language of the lesson content and write ALL questions, options, answers, and explanations in that SAME language. If content is in French, respond in French. If in English, respond in English.

Generate EXACTLY 5 questions in this order and distribution:
- Questions 1â€“2: questionType "TRUE_FALSE", difficulty "EASY"
  * options must be exactly ["Vrai", "Faux"] if content is in French, or ["True", "False"] if in English
  * correctAnswer must be exactly one of the two option values
  * Write ORIGINAL statements in your own words â€” NEVER copy a sentence from the lesson and flip one word
  * Test conceptual understanding, not surface-level recall of wording
- Questions 3â€“4: questionType "MCQ", difficulty "MEDIUM"
  * Exactly 4 options â€” all plausible, wrong ones are related to the topic and wrong for a specific reason
  * correctAnswer must exactly match one of the 4 options
  * Require reasoning: compare concepts, identify correct use case/context, apply concept to scenario, or identify a consequence/implication
  * NEVER ask a question whose answer is a single isolated fact that can be read directly from one sentence
- Question 5: questionType "FILL_BLANK", difficulty "HARD"
  * options must be an empty array []
  * question text MUST contain "____" as the blank placeholder
  * correctAnswer is the exact word or short phrase that fills the blank
  * Target a key technical term that requires genuine understanding â€” NOT a term that appears verbatim in an almost identical sentence in the lesson
  * Construct the sentence around the concept, not around a sentence already in the lesson

Return ONLY a valid JSON array of exactly 5 question objects. No markdown, no code fences, no text outside the JSON array.

Each question object must have these fields:
- "questionType": "TRUE_FALSE" | "MCQ" | "FILL_BLANK"
- "question": the question or statement text (string)
- "options": array of strings (2 for TRUE_FALSE, 4 for MCQ, empty [] for FILL_BLANK)
- "correctAnswer": the correct answer string (must exactly match one option for TRUE_FALSE and MCQ)
- "explanation": why this answer is correct (and for MCQ why the other options are wrong)
- "difficulty": "EASY" | "MEDIUM" | "HARD"

Lesson content:
{lesson_content}"""

    response = client.chat.completions.create(
        model=MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.9,
        max_tokens=4096,
    )
    raw     = response.choices[0].message.content.strip()
    cleaned = clean_json_response(raw)
    return json.loads(cleaned)


class GenerateQuizRequest(BaseModel):
    lessonContent: str
    previousQuestions: list[str] = []


class GenerateExamRequest(BaseModel):
    lessonContents: list[str]
    courseTitle: str


class GenerateLessonRecapRequest(BaseModel):
    lessonId: int
    lessonNumber: int
    lessonTitle: str
    flashcardTerms: list[str]
    lessonSummary: str
    estimatedReadTime: int
    courseTitle: str
    language: str  # "fr", "en", or "ar"


# â”€â”€â”€ On-demand Exam Generation â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
def generate_exam_questions(lesson_contents: list[str], course_title: str) -> list:
    """
    Generate 75 exam questions via 4 separate Groq calls:
      Section 1 â€” 30 TRUE_FALSE  EASY      (1 pt each)   â†’ 1 call
      Section 2 â€” 30 MCQ         MEDIUM    (2 pts each)  â†’ 2 calls of 15
      Section 3 â€” 15 FILL_BLANK  HARD      (3 pts each)  â†’ 1 call
    Returns a flat list of â‰¥75 question dicts.
    """
    combined_content = "\n\n---\n\n".join(lesson_contents)
    content_preview = combined_content[:10000]  # keep well within token budget

    def _call(section_prompt: str) -> list:
        response = client.chat.completions.create(
            model=MODEL,
            messages=[{"role": "user", "content": section_prompt}],
            temperature=0.7,
            max_tokens=6000,
        )
        raw = response.choices[0].message.content.strip()
        cleaned = clean_json_response(raw)
        # Try normal parse first
        try:
            return json.loads(cleaned)
        except json.JSONDecodeError:
            # Attempt to recover a partial array by truncating at the last complete object
            last_brace = cleaned.rfind("},")
            if last_brace != -1:
                partial = cleaned[:last_brace + 1] + "]"
                try:
                    return json.loads(partial)
                except json.JSONDecodeError:
                    pass
            raise

    base_rules = f"""LANGUAGE RULE: Detect the language of the content and write ALL questions, answers, and explanations in that SAME language.
QUALITY RULES:
- NEVER copy a sentence directly from the content â€” rephrase and test understanding
- Test conceptual understanding, real-world application, or relationships between concepts
- No duplicate questions
- Every correctAnswer must be unambiguous and defensible

Course title: {course_title}
Course content (all lessons combined):
{content_preview}"""

    # â”€â”€ Section 1: TRUE_FALSE (30 questions, 1 call) â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    tf_prompt = f"""You are an expert exam question writer for a university e-learning platform.

{base_rules}

Generate EXACTLY 30 TRUE_FALSE questions, difficulty EASY, sectionNumber 1.

Rules:
- option1 must be "Vrai" and option2 must be "Faux" if content is in French, OR "True" and "False" if in English
- option3 and option4 must be null
- correctAnswer must exactly match option1 or option2
- Write ORIGINAL statements that test conceptual understanding

Return ONLY a valid JSON array of exactly 30 objects with these fields:
{{"questionType":"TRUE_FALSE","questionText":"...","option1":"Vrai","option2":"Faux","option3":null,"option4":null,"correctAnswer":"Vrai","explanation":"...","difficulty":"EASY","sectionNumber":1,"pointsWorth":1}}

No markdown, no code fences, ONLY the JSON array."""

    # â”€â”€ Section 2a: MCQ first 15 â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    mcq_prompt_a = f"""You are an expert exam question writer for a university e-learning platform.

{base_rules}

Generate EXACTLY 15 MCQ questions, difficulty MEDIUM, sectionNumber 2.

Rules:
- Exactly 4 options, all plausible
- correctAnswer must exactly match one of the 4 option values
- Test reasoning: comparisons, application, consequences â€” not isolated facts
- Keep explanations concise (1â€“2 sentences)

Return ONLY a valid JSON array of exactly 15 objects with these fields:
{{"questionType":"MCQ","questionText":"...","option1":"...","option2":"...","option3":"...","option4":"...","correctAnswer":"...","explanation":"...","difficulty":"MEDIUM","sectionNumber":2,"pointsWorth":2}}

No markdown, no code fences, ONLY the JSON array."""

    # â”€â”€ Section 2b: MCQ next 15 â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    mcq_prompt_b = f"""You are an expert exam question writer for a university e-learning platform.

{base_rules}

Generate EXACTLY 15 MORE MCQ questions, difficulty MEDIUM, sectionNumber 2.
These must be DIFFERENT from any questions already generated â€” cover different concepts and aspects of the course.

Rules:
- Exactly 4 options, all plausible
- correctAnswer must exactly match one of the 4 option values
- Test reasoning: comparisons, application, consequences â€” not isolated facts
- Keep explanations concise (1â€“2 sentences)

Return ONLY a valid JSON array of exactly 15 objects with these fields:
{{"questionType":"MCQ","questionText":"...","option1":"...","option2":"...","option3":"...","option4":"...","correctAnswer":"...","explanation":"...","difficulty":"MEDIUM","sectionNumber":2,"pointsWorth":2}}

No markdown, no code fences, ONLY the JSON array."""

    # â”€â”€ Section 3: FILL_BLANK (15 questions, 1 call) â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    fb_prompt = f"""You are an expert exam question writer for a university e-learning platform.

{base_rules}

Generate EXACTLY 15 FILL_BLANK questions, difficulty HARD, sectionNumber 3.

Rules:
- The questionText MUST contain "____" as the blank placeholder
- option1, option2, option3, option4 must ALL be null
- correctAnswer is the exact word or short phrase that fills the blank
- Target key technical terms that require deep understanding
- Keep explanations concise (1â€“2 sentences)

Return ONLY a valid JSON array of exactly 15 objects with these fields:
{{"questionType":"FILL_BLANK","questionText":"The ____ is responsible for...","option1":null,"option2":null,"option3":null,"option4":null,"correctAnswer":"exact answer","explanation":"...","difficulty":"HARD","sectionNumber":3,"pointsWorth":3}}

No markdown, no code fences, ONLY the JSON array."""

    def _generate_all() -> list:
        print("[generate_exam_questions] Generating Section 1 (TRUE_FALSE, 30)...")
        section1 = _call(tf_prompt)

        print("[generate_exam_questions] Generating Section 2a (MCQ, 15)...")
        section2a = _call(mcq_prompt_a)

        print("[generate_exam_questions] Generating Section 2b (MCQ, 15)...")
        section2b = _call(mcq_prompt_b)

        print("[generate_exam_questions] Generating Section 3 (FILL_BLANK, 15)...")
        section3 = _call(fb_prompt)

        all_q = section1 + section2a + section2b + section3
        print(f"[generate_exam_questions] Total questions generated: {len(all_q)}")
        return all_q

    try:
        return _generate_all()
    except Exception as e:
        print(f"[generate_exam_questions] First attempt failed: {e}, retrying...")
        try:
            return _generate_all()
        except Exception as retry_err:
            raise RuntimeError(f"Exam generation failed after retry: {retry_err}")


# â”€â”€â”€ Endpoints â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
@app.get("/health")
def health_check():
    return {"status": "ok", "service": "learnai-ai-service", "model": MODEL}


@app.post("/api/generate-quiz")
async def generate_quiz(request: GenerateQuizRequest):
    """Generate 5 fresh quiz questions on demand from lesson content."""
    if not request.lessonContent or len(request.lessonContent.strip()) < 50:
        raise HTTPException(status_code=400, detail="lessonContent is too short to generate questions.")
    try:
        questions = generate_quiz_questions(request.lessonContent, request.previousQuestions)
    except (json.JSONDecodeError, Exception):
        # Retry once on parse failure
        try:
            questions = generate_quiz_questions(request.lessonContent, request.previousQuestions)
        except Exception as retry_err:
            raise HTTPException(status_code=502, detail=f"Quiz generation failed: {str(retry_err)}")
    return questions


@app.post("/api/generate-exam")
async def generate_exam(request: GenerateExamRequest):
    """Generate 75 exam questions (30 TF + 30 MCQ + 15 FB) from all lesson contents."""
    if not request.lessonContents or all(len(c.strip()) < 50 for c in request.lessonContents):
        raise HTTPException(status_code=400, detail="lessonContents is empty or too short.")
    try:
        questions = generate_exam_questions(request.lessonContents, request.courseTitle)
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Exam generation failed: {str(e)}")
    return questions


@app.post("/api/generate-lesson-recap")
async def generate_lesson_recap_endpoint(request: GenerateLessonRecapRequest):
    """Generate (or return cached) a recap video for a single lesson."""
    language_code = normalize_video_language(request.language)
    safe = re.sub(r"[^\w]", "_", request.lessonTitle.encode("ascii", "ignore").decode())[:30]
    expected_path = os.path.abspath(
        os.path.join("uploads", "recap-videos", f"lesson_{request.lessonNumber}_{language_code}_{safe}.mp4")
    )
    # Return cached file if it already exists
    if os.path.exists(expected_path):
        print(f"[generate-lesson-recap] cache hit: {expected_path}")
        return _recap_response(expected_path)
    # Generate
    try:
        result = generate_recap_video(
            lesson_number      = request.lessonNumber,
            lesson_title       = request.lessonTitle,
            flashcard_terms    = request.flashcardTerms,
            lesson_summary     = request.lessonSummary,
            estimated_read_time= request.estimatedReadTime,
            course_title       = request.courseTitle,
            language           = language_code,
        )
        return result or {"recapVideoPath": ""}
    except HTTPException:
        raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=502, detail=f"Video generation failed: {str(e)}")


# â”€â”€â”€ Certificate generation â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

class CertificateRequest(BaseModel):
    certificateUuid: str
    studentName: str
    courseTitle: str
    score: int
    issuedAt: str


def generate_certificate_html(
    certificate_uuid: str,
    student_name: str,
    course_title: str,
    score: int,
    issued_at: str,
) -> str:
    """Return a self-contained A4-landscape HTML string (kept for reference)."""
    try:
        from datetime import datetime
        dt = datetime.fromisoformat(issued_at.replace("Z", "+00:00"))
        date_str = dt.strftime("%B %d, %Y")
    except Exception:
        date_str = issued_at[:10]

    return f"""<!DOCTYPE html><html><body>
    <h1>Certificate of Achievement</h1>
    <p>{student_name} â€” {course_title} â€” Score: {score}% â€” {date_str}</p>
    </body></html>"""


def generate_certificate_pdf(
    certificate_uuid: str,
    student_name: str,
    course_title: str,
    score: int,
    issued_at: str,
) -> str:

    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib import colors
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import os, re
    from datetime import datetime

    # -- Output path -------------------------------
    output_path = f"/tmp/certificate_{certificate_uuid}.pdf"
    page_width, page_height = landscape(A4)
    c = canvas.Canvas(output_path,
                      pagesize=landscape(A4))

    # -- Parse date --------------------------------
    try:
        dt = datetime.fromisoformat(
            issued_at.replace("Z", "+00:00"))
        date_str = dt.strftime("%B %d, %Y")
    except Exception:
        date_str = str(issued_at)[:10]

    # -- Color palette -----------------------------
    dark_bg      = colors.HexColor("#0D0F1A")
    dark_card    = colors.HexColor("#12152A")
    purple       = colors.HexColor("#6366F1")
    purple_light = colors.HexColor("#818CF8")
    gold         = colors.HexColor("#F59E0B")
    gold_light   = colors.HexColor("#FCD34D")
    white        = colors.HexColor("#F8FAFC")
    grey         = colors.HexColor("#94A3B8")
    dark_grey    = colors.HexColor("#334155")

    # -- Background --------------------------------
    c.setFillColor(dark_bg)
    c.rect(0, 0, page_width, page_height,
           fill=1, stroke=0)

    # -- Outer decorative border -------------------
    border_margin = 12 * mm
    c.setStrokeColor(purple)
    c.setLineWidth(2)
    c.rect(border_margin, border_margin,
           page_width - 2 * border_margin,
           page_height - 2 * border_margin,
           fill=0, stroke=1)

    # Inner gold border
    inner_margin = 16 * mm
    c.setStrokeColor(gold)
    c.setLineWidth(0.5)
    c.rect(inner_margin, inner_margin,
           page_width - 2 * inner_margin,
           page_height - 2 * inner_margin,
           fill=0, stroke=1)

    # -- Corner decorations ------------------------
    corner_size = 8 * mm
    corners = [
        (inner_margin, inner_margin),
        (page_width - inner_margin, inner_margin),
        (inner_margin, page_height - inner_margin),
        (page_width - inner_margin,
         page_height - inner_margin),
    ]
    c.setStrokeColor(gold)
    c.setLineWidth(1.5)
    for cx, cy in corners:
        c.circle(cx, cy, 3 * mm, fill=0, stroke=1)

    # -- Header - Platform name --------------------
    c.setFillColor(purple_light)
    c.setFont("Helvetica-Bold", 11)
    platform = "?  L E A R N A I  P L A T F O R M  ?"
    c.drawCentredString(
        page_width / 2,
        page_height - 28 * mm,
        platform)

    # Divider line under platform name
    c.setStrokeColor(purple)
    c.setLineWidth(0.8)
    c.line(inner_margin + 10 * mm,
           page_height - 32 * mm,
           page_width - inner_margin - 10 * mm,
           page_height - 32 * mm)

    # -- Main title --------------------------------
    c.setFillColor(white)
    c.setFont("Helvetica-Bold", 28)
    c.drawCentredString(
        page_width / 2,
        page_height - 48 * mm,
        "CERTIFICATE OF ACHIEVEMENT")

    # Gold underline for title
    title_width = 180 * mm
    c.setStrokeColor(gold)
    c.setLineWidth(1.2)
    c.line(page_width / 2 - title_width / 2,
           page_height - 51 * mm,
           page_width / 2 + title_width / 2,
           page_height - 51 * mm)

    # -- Subtitle ----------------------------------
    c.setFillColor(grey)
    c.setFont("Helvetica", 11)
    c.drawCentredString(
        page_width / 2,
        page_height - 62 * mm,
        "This is to certify that")

    # -- Student name ------------------------------
    c.setFillColor(gold_light)
    c.setFont("Helvetica-Bold", 32)
    c.drawCentredString(
        page_width / 2,
        page_height - 78 * mm,
        student_name)

    # Name underline
    name_width = c.stringWidth(
        student_name, "Helvetica-Bold", 32)
    c.setStrokeColor(gold)
    c.setLineWidth(1)
    c.line(page_width / 2 - name_width / 2 - 5,
           page_height - 81 * mm,
           page_width / 2 + name_width / 2 + 5,
           page_height - 81 * mm)

    # -- Course label ------------------------------
    c.setFillColor(grey)
    c.setFont("Helvetica", 11)
    c.drawCentredString(
        page_width / 2,
        page_height - 91 * mm,
        "has successfully completed the course")

    # -- Course title ------------------------------
    max_chars = 55
    if len(course_title) > max_chars:
        display_title = (
            '"' + course_title[:max_chars - 3]
            + '..."')
    else:
        display_title = '"' + course_title + '"'

    c.setFillColor(purple_light)
    c.setFont("Helvetica-Bold", 16)
    c.drawCentredString(
        page_width / 2,
        page_height - 103 * mm,
        display_title)

    # -- Score badge -------------------------------
    badge_w = 40 * mm
    badge_h = 10 * mm
    badge_x = page_width / 2 - badge_w / 2
    badge_y = page_height - 118 * mm

    c.setFillColor(purple)
    c.roundRect(badge_x, badge_y,
                badge_w, badge_h,
                3 * mm, fill=1, stroke=0)

    c.setFillColor(white)
    c.setFont("Helvetica-Bold", 12)
    c.drawCentredString(
        page_width / 2,
        badge_y + 3 * mm,
        f"Score: {score}%")

    # -- Bottom divider ----------------------------
    divider_y = 38 * mm
    c.setStrokeColor(dark_grey)
    c.setLineWidth(0.5)
    c.line(inner_margin + 10 * mm,
           divider_y,
           page_width - inner_margin - 10 * mm,
           divider_y)

    # -- Footer - 3 columns ------------------------
    col1_x = page_width * 0.2
    col2_x = page_width * 0.5
    col3_x = page_width * 0.8
    footer_label_y = 28 * mm
    footer_value_y = 22 * mm

    # Column labels
    c.setFillColor(grey)
    c.setFont("Helvetica", 8)
    for label, x in [
        ("DATE ISSUED", col1_x),
        ("VERIFICATION ID", col2_x),
        ("PLATFORM SEAL", col3_x),
    ]:
        c.drawCentredString(x, footer_label_y,
                            label)

    # Date value
    c.setFillColor(white)
    c.setFont("Helvetica-Bold", 10)
    c.drawCentredString(
        col1_x, footer_value_y, date_str)

    # Full UUID value (split into two lines
    # if needed)
    uuid_line1 = certificate_uuid[:36]
    c.setFillColor(purple_light)
    c.setFont("Helvetica", 8)
    c.drawCentredString(
        col2_x, footer_value_y + 3 * mm,
        uuid_line1)

    # Platform seal - decorative circle
    seal_x = col3_x
    seal_y = footer_value_y + 2 * mm
    c.setFillColor(purple)
    c.circle(seal_x, seal_y,
             7 * mm, fill=1, stroke=0)
    c.setStrokeColor(gold)
    c.setLineWidth(1)
    c.circle(seal_x, seal_y,
             7 * mm, fill=0, stroke=1)
    c.setFillColor(white)
    c.setFont("Helvetica-Bold", 7)
    c.drawCentredString(seal_x, seal_y + 1, "?")
    c.setFont("Helvetica", 5)
    c.drawCentredString(
        seal_x, seal_y - 3, "VERIFIED")

    # -- Subtle background glow effect -------------
    # Purple glow top center
    c.setFillColor(
        colors.HexColor("#6366F1"))
    c.setFillAlpha(0.04)
    c.circle(page_width / 2,
             page_height - 10 * mm,
             80 * mm, fill=1, stroke=0)
    c.setFillAlpha(1)

    c.save()
    return output_path


@app.post("/api/generate-certificate")
async def generate_certificate_endpoint(req: CertificateRequest):
    """Generate a PDF certificate, return the content as base64, and delete the file."""
    try:
        pdf_path = generate_certificate_pdf(
            certificate_uuid=req.certificateUuid,
            student_name=req.studentName,
            course_title=req.courseTitle,
            score=req.score,
            issued_at=req.issuedAt,
        )
        with open(pdf_path, "rb") as f:
            pdf_bytes = f.read()
        os.remove(pdf_path)
        pdf_base64 = base64.b64encode(pdf_bytes).decode("utf-8")
        return {"pdfContent": pdf_base64, "status": "generated"}
    except HTTPException:
        raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Certificate generation failed: {str(e)}")


@app.post("/process-document")
async def process_document(
    file: UploadFile = File(...),
    fileType: str    = Query(..., description="PDF | DOCX | PPTX | IMAGE"),
):
    # â”€â”€ Validate file type â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    file_type = fileType.upper()
    if file_type not in SUPPORTED_FILE_TYPES:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type '{fileType}'. Supported: {', '.join(SUPPORTED_FILE_TYPES)}",
        )

    # â”€â”€ Read file bytes â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    file_bytes = await file.read()
    if not file_bytes:
        raise HTTPException(status_code=400, detail="Uploaded file is empty or unreadable.")

    # â”€â”€ Extract text â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    print(f"[1/4] Extracting text from {file_type} file: {file.filename}")
    try:
        text = extract_text(file_bytes, file_type, file.filename or "")
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=422, detail=f"Text extraction failed: {str(e)}")

    if len(text) < MIN_TEXT_LENGTH:
        raise HTTPException(
            status_code=422,
            detail="Document does not contain enough readable content to generate a course (minimum 50 characters).",
        )
    print(f"    Extracted {len(text)} characters of text.")

    # â”€â”€ Chunk text â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    chunks = split_into_chunks(text)
    print(f"[2/4] Split into {len(chunks)} chunk(s) (max {MAX_CHUNKS}, ~{CHUNK_WORD_SIZE} words each).")

    # â”€â”€ Category detection â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    print("[3/4] Detecting document category...")
    try:
        category = detect_category(text)
    except Exception as e:
        category = "GÃ©nÃ©ral"
    print(f"    Category: \"{category}\"")

    # â”€â”€ Course title â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    print("[4/4] Generating course title...")
    try:
        course_title = generate_course_title(text)
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Groq API error (title): {str(e)}")
    print(f"    Title: \"{course_title}\"")

    # â”€â”€ Course description â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    try:
        course_description = generate_course_description(text)
    except Exception:
        course_description = ""

    # â”€â”€ Lessons â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    lessons = []
    previous_titles: list[str] = []
    print(f"[5/5] Generating {len(chunks)} lesson(s)...")
    for idx, chunk in enumerate(chunks, start=1):
        print(f"    Lesson {idx}/{len(chunks)}...")
        try:
            lesson = generate_lesson(chunk, idx, previous_titles)
        except (json.JSONDecodeError, Exception):
            print(f"    Lesson {idx} parse failed, retrying...")
            try:
                lesson = generate_lesson(chunk, idx, previous_titles)
            except Exception as retry_err:
                raise HTTPException(
                    status_code=502,
                    detail=f"Groq API error on lesson {idx}: {str(retry_err)}",
                )
        lessons.append(lesson)
        previous_titles.append(lesson.get("title", f"Lesson {idx}"))

    result = {
        "courseTitle":       course_title,
        "courseDescription": course_description,
        "category":          category,
        "totalLessons":      len(lessons),
        "lessons":           lessons,
    }
    print(f"Done. Course \"{course_title}\" â€” category: \"{category}\" â€” {len(lessons)} lesson(s).")
    return result


# â”€â”€â”€ Entry point â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)


# â”€â”€â”€ Chatbot â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

class ChatMessage(BaseModel):
    role: str
    content: str


class ChatRequest(BaseModel):
    messages: list[ChatMessage]
    mode: str = "visitor"
    context: Optional[str] = None


def groq_chat_multi(system_prompt: str, messages: list[ChatMessage]) -> str:
    """Multi-turn chat completion call to Groq with a system prompt."""
    built_messages = [{"role": "system", "content": system_prompt}] + [
        {"role": m.role, "content": m.content} for m in messages
    ]
    response = client.chat.completions.create(
        model=MODEL,
        messages=built_messages,
        max_tokens=1024,
    )
    return response.choices[0].message.content.strip()


@app.post("/api/chat")
async def chat_endpoint(request: ChatRequest):
    """Multi-turn AI chatbot. Supports visitor (platform FAQ) and student (lesson help) modes."""
    if not request.messages:
        raise HTTPException(status_code=400, detail="messages array must not be empty.")

    if request.mode == "student":
        if request.context:
            system_prompt = f"""You are a helpful learning 
assistant for the LearnAI platform. The student is 
currently studying the following lesson content:

{request.context}

Answer questions based on this lesson content. 
Be concise, clear, and educational."""
        else:
            system_prompt = """You are a helpful learning 
assistant for the LearnAI platform. Help the student 
with their learning questions."""
    else:
        # visitor mode (default)
        system_prompt = (
            "You are the LearnAI assistant. "
            "LearnAI is an AI-powered personalized learning platform. "
            "Students upload their academic documents (PDF, DOCX) and the platform automatically generates "
            "a complete personalized course including lessons, quizzes, flashcards, and a final exam. "
            "Students can earn certificates after passing the final exam. "
            "The platform uses Llama 3.3 70B for content generation, ElevenLabs for lesson recap narration, "
            "and includes an anti-cheating system during exams. "
            "Answer questions about the platform clearly and concisely. "
            "If asked something unrelated to the platform, politely redirect the conversation back to LearnAI."
        )

    try:
        reply = groq_chat_multi(system_prompt, request.messages)
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Chat generation failed: {str(e)}")

    return {"reply": reply}


